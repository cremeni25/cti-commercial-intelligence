from __future__ import annotations

from io import BytesIO
from typing import Any

from docx import Document
from openpyxl import Workbook
from openpyxl.styles import Font
from pptx import Presentation
from pptx.util import Inches, Pt


def _artefatos(metadados: dict[str, Any]) -> list[dict[str, Any]]:
    return [item for item in (metadados.get("artefatos") or []) if isinstance(item, dict)]


def _fontes(metadados: dict[str, Any]) -> list[dict[str, Any]]:
    return [item for item in (metadados.get("fontes") or []) if isinstance(item, dict)]


def _serie(metadados: dict[str, Any]) -> list[dict[str, Any]]:
    for item in _artefatos(metadados):
        dados = item.get("dados")
        if isinstance(dados, list) and dados:
            return [linha for linha in dados if isinstance(linha, dict)]
    return []


def gerar_xlsx_resposta(conteudo: str, metadados: dict[str, Any]) -> bytes:
    wb = Workbook()
    ws = wb.active
    ws.title = "Análise CTI"
    ws["A1"] = "IA Comercial CTI — Refrigeração de Transporte e Cadeia Fria"
    ws["A1"].font = Font(bold=True, size=14)
    ws["A3"] = "Resposta"
    ws["A3"].font = Font(bold=True)
    for indice, linha in enumerate(str(conteudo or "").splitlines() or [""], start=4):
        ws.cell(indice, 1, linha)
    ws.column_dimensions["A"].width = 90

    serie = _serie(metadados)
    if serie:
        dados = wb.create_sheet("Dados")
        dados.append(["Label", "Valor", "Unidade"])
        for celula in dados[1]:
            celula.font = Font(bold=True)
        for item in serie:
            dados.append([item.get("label"), item.get("valor"), item.get("unidade")])
        dados.column_dimensions["A"].width = 45
        dados.column_dimensions["B"].width = 18
        dados.column_dimensions["C"].width = 18

    fontes = _fontes(metadados)
    if fontes:
        ws_fontes = wb.create_sheet("Fontes")
        ws_fontes.append(["Tipo", "Descrição", "URL"])
        for celula in ws_fontes[1]:
            celula.font = Font(bold=True)
        for fonte in fontes:
            ws_fontes.append([fonte.get("tipo"), fonte.get("descricao"), fonte.get("url")])
        ws_fontes.column_dimensions["A"].width = 18
        ws_fontes.column_dimensions["B"].width = 60
        ws_fontes.column_dimensions["C"].width = 80

    saida = BytesIO()
    wb.save(saida)
    return saida.getvalue()


def gerar_docx_resposta(conteudo: str, metadados: dict[str, Any]) -> bytes:
    doc = Document()
    doc.add_heading("IA Comercial CTI", level=0)
    doc.add_paragraph("Refrigeração de Transporte e Cadeia Fria")
    for bloco in str(conteudo or "").split("\n\n"):
        texto = bloco.strip()
        if texto:
            doc.add_paragraph(texto)

    fontes = _fontes(metadados)
    if fontes:
        doc.add_heading("Fontes", level=1)
        for fonte in fontes:
            descricao = str(fonte.get("descricao") or fonte.get("tipo") or "Fonte")
            url = str(fonte.get("url") or "")
            doc.add_paragraph(f"{descricao}{' — ' + url if url else ''}", style="List Bullet")

    saida = BytesIO()
    doc.save(saida)
    return saida.getvalue()


def _blocos_apresentacao(conteudo: str, limite: int = 6) -> list[str]:
    paragrafos = [p.strip() for p in str(conteudo or "").split("\n\n") if p.strip()]
    if not paragrafos:
        return ["Sem conteúdo textual disponível."]
    blocos: list[str] = []
    atual = ""
    for paragrafo in paragrafos:
        candidato = f"{atual}\n\n{paragrafo}".strip()
        if atual and len(candidato) > 1200:
            blocos.append(atual)
            atual = paragrafo
        else:
            atual = candidato
        if len(blocos) >= limite - 1:
            break
    if atual and len(blocos) < limite:
        blocos.append(atual)
    return blocos


def gerar_pptx_resposta(conteudo: str, metadados: dict[str, Any]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    capa = prs.slides.add_slide(prs.slide_layouts[0])
    capa.shapes.title.text = "IA Comercial CTI"
    capa.placeholders[1].text = "Inteligência para Refrigeração de Transporte e Cadeia Fria"

    for numero, bloco in enumerate(_blocos_apresentacao(conteudo), start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Análise {numero}"
        caixa = slide.placeholders[1].text_frame
        caixa.clear()
        p = caixa.paragraphs[0]
        p.text = bloco
        p.font.size = Pt(18)

    fontes = _fontes(metadados)
    if fontes:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Fontes"
        caixa = slide.placeholders[1].text_frame
        caixa.clear()
        for indice, fonte in enumerate(fontes[:12]):
            p = caixa.paragraphs[0] if indice == 0 else caixa.add_paragraph()
            p.text = str(fonte.get("descricao") or fonte.get("tipo") or "Fonte")
            if fonte.get("url"):
                p.text += f" — {fonte['url']}"
            p.font.size = Pt(14)

    saida = BytesIO()
    prs.save(saida)
    return saida.getvalue()
