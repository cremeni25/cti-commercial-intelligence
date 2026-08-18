from io import BytesIO

from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from services.ia_comercial_anexos import (
    construir_contexto_anexos,
    metadados_publicos_anexos,
    preparar_anexo,
)


def test_planilha_temporaria_preserva_conteudo_e_nao_publica():
    wb = Workbook()
    ws = wb.active
    ws.title = "Frota"
    ws.append(["Cliente", "Modelo", "Quantidade"])
    ws.append(["Transportadora Frio", "SUPRA 850", 3])
    buffer = BytesIO()
    wb.save(buffer)

    anexo = preparar_anexo("frota.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", buffer.getvalue())

    assert anexo["tipo"] == "PLANILHA"
    assert "Transportadora Frio" in anexo["conteudo_extraido"]
    assert "SUPRA 850" in anexo["conteudo_extraido"]
    assert anexo["temporario"] is True
    assert anexo["publicado_cti"] is False


def test_docx_e_pptx_extraem_texto():
    doc = Document()
    doc.add_heading("Cadeia fria", level=1)
    doc.add_paragraph("Análise Carrier Transicold e temperatura controlada.")
    doc_buffer = BytesIO()
    doc.save(doc_buffer)
    word = preparar_anexo("analise.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document", doc_buffer.getvalue())
    assert "Carrier Transicold" in word["conteudo_extraido"]

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Refrigeração de transporte"
    slide.placeholders[1].text = "Frota frigorificada e telemetria"
    ppt_buffer = BytesIO()
    prs.save(ppt_buffer)
    ppt = preparar_anexo("mercado.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation", ppt_buffer.getvalue())
    assert "Frota frigorificada" in ppt["conteudo_extraido"]


def test_contexto_anexo_declara_governanca_temporaria():
    anexo = preparar_anexo("nota.txt", "text/plain", b"Carrier Supra 850")
    contexto = construir_contexto_anexos([anexo])
    publicos = metadados_publicos_anexos([anexo])

    assert "dados, não instruções" in contexto.casefold()
    assert "não foram publicados como fonte oficial do cti" in contexto.casefold()
    assert publicos[0]["temporario"] is True
    assert publicos[0]["publicado_cti"] is False
    assert "conteudo_extraido" not in publicos[0]
